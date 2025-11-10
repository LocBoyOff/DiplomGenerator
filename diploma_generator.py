import os
import json
import threading
import queue
import time
import psutil
import wx
import wx.grid
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import comtypes.client
from datetime import datetime
import img2pdf
import re
import io
from PIL import Image

# --- Backend Functions (Unchanged) ---
def pptx_to_pdf(input_pptx, output_pdf, stop_event):
    powerpoint = None
    try:
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = 1
        deck = powerpoint.Presentations.Open(input_pptx)
        if stop_event.is_set():
            deck.Close()
            return False
        temp_jpg = input_pptx.replace(".pptx", ".jpg")
        deck.Slides[1].Export(temp_jpg, "JPG", 3508, 2480)
        deck.Close()
        
        if stop_event.is_set():
            os.remove(temp_jpg) if os.path.exists(temp_jpg) else None
            return False
        
        a4inpt = (img2pdf.mm_to_pt(297), img2pdf.mm_to_pt(210))
        layout = img2pdf.get_layout_fun(a4inpt)
        with open(output_pdf, "wb") as f:
            f.write(img2pdf.convert(temp_jpg, layout_fun=layout))
        
        os.remove(temp_jpg)
        return True
    except Exception as e:
        raise Exception(f"Ошибка конвертации: {e}")
    finally:
        if powerpoint:
            try:
                powerpoint.Quit()
            except:
                for proc in psutil.process_iter(['name']):
                    if proc.info['name'].lower() == 'powerpnt.exe':
                        proc.terminate()
                        try:
                            proc.wait(timeout=3)
                        except psutil.TimeoutExpired:
                            proc.kill()

def replace_text(shape, placeholder, value, font_settings=None):
    if shape.has_text_frame:
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            original_text = paragraph.text
            if placeholder in original_text:
                runs = paragraph.runs
                if runs:
                    full_text = "".join(run.text for run in runs)
                    new_text = full_text.replace(placeholder, value)
                    paragraph.clear()
                    new_run = paragraph.add_run()
                    new_run.text = new_text
                    original_font = runs[0].font
                    if font_settings and font_settings["use_custom"]:
                        new_run.font.name = font_settings["name"]
                        new_run.font.size = Pt(font_settings["size"])
                        new_run.font.bold = font_settings.get("bold", False)
                    else:
                        new_run.font.size = original_font.size
                        new_run.font.bold = original_font.bold
                        new_run.font.name = original_font.name
                    new_run.font.color.rgb = RGBColor(127, 127, 127)
                    paragraph.alignment = paragraph.alignment or PP_ALIGN.CENTER
                else:
                    paragraph.text = original_text.replace(placeholder, value)
                    for run in paragraph.runs:
                        if font_settings and font_settings["use_custom"]:
                            run.font.name = font_settings["name"]
                            run.font.size = Pt(font_settings["size"])
                            run.font.bold = font_settings.get("bold", False)
                        run.font.color.rgb = RGBColor(127, 127, 127)
                    paragraph.alignment = PP_ALIGN.CENTER

def generate_diplomas(excel_path, ppt_template, output_dir, column_mapping, error_handling, default_values, font_settings, sort_column, enable_sorting, log_queue, progress_queue, eta_queue, stop_event):
    wb = load_workbook(excel_path)
    ws = wb.active
    participants = []
    skipped_rows = []
    
    headers = [cell.value or f"Столбец {chr(65+i)}" for i, cell in enumerate(ws[1])]
    for row_idx, row in enumerate(ws.iter_rows(min_row=2, values_only=True), start=2):
        if stop_event.is_set():
            log_queue.put("Генерация прервана")
            return False
        participant = {}
        valid = True
        for placeholder, col_name in column_mapping.items():
            col_idx = headers.index(col_name) if col_name in headers else ord(col_name) - ord('A')
            value = row[col_idx]
            if value is None:
                if error_handling == "skip":
                    skipped_rows.append(f"Строка {row_idx}: пустое поле {placeholder} ({col_name})")
                    valid = False
                    break
                elif error_handling == "default":
                    value = default_values.get(placeholder, "Не указано")
                else:
                    log_queue.put(f"Ошибка: пустое поле {placeholder} ({col_name}) в строке {row_idx}")
                    return False
            if placeholder == "DATE" and isinstance(value, datetime):
                value = value.strftime("%d.%m.%Y")
            elif placeholder == "DATE" and isinstance(value, str):
                for fmt in ["%Y-%m-%d", "%d.%m.%Y", "%d/%m/%Y", "%Y/%m/%d"]:
                    try:
                        value = datetime.strptime(value, fmt).strftime("%d.%m.%Y")
                        break
                    except ValueError:
                        pass
            participant[placeholder] = str(value)
        if valid:
            participants.append(participant)
    
    if skipped_rows:
        log_queue.put(f"Пропущены строки: {len(skipped_rows)}. Подробности: {'; '.join(skipped_rows)}")
    
    os.makedirs(output_dir, exist_ok=True)
    
    total = len(participants)
    processing_times = []
    for idx, participant in enumerate(participants, 1):
        if stop_event.is_set():
            log_queue.put("Генерация прервана")
            return False
        start_time = time.time()
        prs = Presentation(ppt_template)
        slide = prs.slides[0]
        
        for shape in slide.shapes:
            for placeholder in participant:
                replace_text(shape, "{" + placeholder + "}", participant[placeholder], font_settings)
        
        safe_name = re.sub(r'[\\/*?:"<>|]', "_", participant.get("NAME", "unknown"))
        pdf_name = f"{safe_name}.pdf"
        
        if enable_sorting and sort_column:
            safe_sort_value = re.sub(r'[\\/*?:"<>|]', "_", participant.get(sort_column, "unknown"))
            subdir = os.path.join(output_dir, safe_sort_value)
            os.makedirs(subdir, exist_ok=True)
            pdf_path = os.path.join(subdir, pdf_name)
        else:
            pdf_path = os.path.join(output_dir, pdf_name)
        
        temp_pptx = os.path.abspath(f"temp_{safe_name}.pptx")
        prs.save(temp_pptx)
        try:
            if not pptx_to_pdf(temp_pptx, pdf_path, stop_event):
                return False
            log_queue.put(f"Сгенерирован диплом: {pdf_name}")
        except Exception as e:
            log_queue.put(str(e))
            return False
        finally:
            os.remove(temp_pptx) if os.path.exists(temp_pptx) else None
        processing_time = time.time() - start_time
        processing_times.append(processing_time)
        
        progress_queue.put(idx / total * 100)
        if processing_times:
            avg_time = sum(processing_times) / len(processing_times)
            remaining_records = total - idx
            eta_seconds = int(avg_time * remaining_records)
            eta_str = f"{eta_seconds // 60:02d}:{eta_seconds % 60:02d}"
            eta_queue.put(eta_str)
    
    log_queue.put(f"Дипломы сохранены в: {output_dir}")
    eta_queue.put("00:00")
    return True

# --- GUI Application (wxPython, Updated UI) ---
class DiplomaGeneratorApp(wx.Frame):
    def __init__(self):
        super().__init__(None, title="Генератор дипломов", size=(1000, 400))
        self.SetMinSize((1000, 400))
        self.SetMaxSize((1000, 400))
        
        self.excel_path = ""
        self.pptx_path = ""
        self.output_dir = ""
        self.column_mapping = {}
        self.placeholders = []
        self.error_handling = "stop"
        self.default_values = {}
        self.font_settings = {"use_custom": False}
        self.sort_column = ""
        self.enable_sorting = True
        self.stop_event = threading.Event()
        self.log_queue = queue.Queue()
        self.progress_queue = queue.Queue()
        self.eta_queue = queue.Queue()
        self.generation_thread = None
        
        self.setup_ui()
        self.load_config()
        self.Bind(wx.EVT_CLOSE, self.on_closing)
        
        # Timer for queue checking
        self.timer = wx.Timer(self)
        self.Bind(wx.EVT_TIMER, self.check_queues, self.timer)
        self.timer.Start(100)
    
    def setup_ui(self):
        panel = wx.Panel(self)
        main_sizer = wx.BoxSizer(wx.HORIZONTAL)
        
        # Theme
        self.theme = {"bg": wx.Colour(255, 255, 255), "fg": wx.Colour(0, 0, 0), "button_bg": wx.Colour(76, 175, 80)}
        panel.SetBackgroundColour(self.theme["bg"])
        
        # Fonts with Cyrillic support
        self.label_font = wx.Font(12, wx.FONTFAMILY_SWISS, wx.FONTSTYLE_NORMAL, wx.FONTWEIGHT_NORMAL, False, "Arial")
        self.log_font = wx.Font(10, wx.FONTFAMILY_SWISS, wx.FONTSTYLE_NORMAL, wx.FONTWEIGHT_NORMAL, False, "Arial")
        self.button_font = wx.Font(10, wx.FONTFAMILY_SWISS, wx.FONTSTYLE_NORMAL, wx.FONTWEIGHT_NORMAL, False, "Arial")
        
        # Left panel: Inputs and controls
        left_panel = wx.Panel(panel)
        left_sizer = wx.GridBagSizer(vgap=5, hgap=5)
        
        # Excel file
        excel_label = wx.StaticText(left_panel, label="Excel-файл:")
        excel_label.SetFont(self.label_font)
        self.excel_name = wx.StaticText(left_panel, label="Не выбран")
        self.excel_name.SetFont(self.label_font)
        self.excel_path_ctrl = wx.TextCtrl(left_panel, style=wx.TE_READONLY)
        self.excel_path_ctrl.SetToolTip("Полный путь к файлу Excel")
        excel_browse_btn = wx.Button(left_panel, label="Выбрать", size=(100, 30))
        excel_browse_btn.SetBackgroundColour(self.theme["button_bg"])
        excel_browse_btn.SetForegroundColour(wx.Colour(255, 255, 255))
        excel_browse_btn.Bind(wx.EVT_BUTTON, self.choose_excel)
        
        left_sizer.Add(excel_label, pos=(0, 0), flag=wx.ALIGN_CENTER_VERTICAL)
        left_sizer.Add(self.excel_name, pos=(0, 1), flag=wx.ALIGN_CENTER_VERTICAL | wx.RIGHT, border=5)
        left_sizer.Add(self.excel_path_ctrl, pos=(0, 2), flag=wx.EXPAND)
        left_sizer.Add(excel_browse_btn, pos=(0, 3))
        
        # PPTX template
        pptx_label = wx.StaticText(left_panel, label="Шаблон PPTX:")
        pptx_label.SetFont(self.label_font)
        self.pptx_name = wx.StaticText(left_panel, label="Не выбран")
        self.pptx_name.SetFont(self.label_font)
        self.pptx_path_ctrl = wx.TextCtrl(left_panel, style=wx.TE_READONLY)
        self.pptx_path_ctrl.SetToolTip("Полный путь к файлу PPTX")
        pptx_browse_btn = wx.Button(left_panel, label="Выбрать", size=(100, 30))
        pptx_browse_btn.SetBackgroundColour(self.theme["button_bg"])
        pptx_browse_btn.SetForegroundColour(wx.Colour(255, 255, 255))
        pptx_browse_btn.Bind(wx.EVT_BUTTON, self.choose_pptx)
        
        left_sizer.Add(pptx_label, pos=(1, 0), flag=wx.ALIGN_CENTER_VERTICAL)
        left_sizer.Add(self.pptx_name, pos=(1, 1), flag=wx.ALIGN_CENTER_VERTICAL | wx.RIGHT, border=5)
        left_sizer.Add(self.pptx_path_ctrl, pos=(1, 2), flag=wx.EXPAND)
        left_sizer.Add(pptx_browse_btn, pos=(1, 3))
        
        # Output directory
        output_label = wx.StaticText(left_panel, label="Папка:")
        output_label.SetFont(self.label_font)
        self.output_name = wx.StaticText(left_panel, label="Не выбрана")
        self.output_name.SetFont(self.label_font)
        self.output_path_ctrl = wx.TextCtrl(left_panel, style=wx.TE_READONLY)
        self.output_path_ctrl.SetToolTip("Полный путь к папке")
        output_browse_btn = wx.Button(left_panel, label="Выбрать", size=(100, 30))
        output_browse_btn.SetBackgroundColour(self.theme["button_bg"])
        output_browse_btn.SetForegroundColour(wx.Colour(255, 255, 255))
        output_browse_btn.Bind(wx.EVT_BUTTON, self.choose_output)
        
        left_sizer.Add(output_label, pos=(2, 0), flag=wx.ALIGN_CENTER_VERTICAL)
        left_sizer.Add(self.output_name, pos=(2, 1), flag=wx.ALIGN_CENTER_VERTICAL | wx.RIGHT, border=5)
        left_sizer.Add(self.output_path_ctrl, pos=(2, 2), flag=wx.EXPAND)
        left_sizer.Add(output_browse_btn, pos=(2, 3))
        
        # Error handling
        error_label = wx.StaticText(left_panel, label="Ошибки:")
        error_label.SetFont(self.label_font)
        self.error_handling_choice = wx.Choice(left_panel, choices=["Остановить", "Пропустить", "Заполнить по умолчанию"])
        self.error_handling_choice.SetSelection(0)
        self.error_handling_choice.SetToolTip("Как обрабатывать пустые поля в Excel")
        self.error_handling_choice.Bind(wx.EVT_CHOICE, self.update_error_handling)
        
        left_sizer.Add(error_label, pos=(3, 0), flag=wx.ALIGN_CENTER_VERTICAL)
        left_sizer.Add(self.error_handling_choice, pos=(3, 1), span=(1, 2), flag=wx.EXPAND)
        
        # Sorting settings
        sort_label = wx.StaticText(left_panel, label="Сортировка:")
        sort_label.SetFont(self.label_font)
        self.sort_check = wx.CheckBox(left_panel, label="Сортировать в папки")
        self.sort_check.SetValue(self.enable_sorting)
        self.sort_check.SetToolTip("Включить сортировку дипломов в папки")
        self.sort_check.Enable(False)
        self.sort_check.Bind(wx.EVT_CHECKBOX, self.update_sorting)
        self.sort_choice = wx.Choice(left_panel, choices=[])
        self.sort_choice.Enable(False)
        self.sort_choice.SetToolTip("Выберите столбец для сортировки")
        self.sort_choice.Bind(wx.EVT_CHOICE, self.update_sorting)
        
        left_sizer.Add(sort_label, pos=(4, 0), flag=wx.ALIGN_CENTER_VERTICAL)
        left_sizer.Add(self.sort_check, pos=(4, 1))
        left_sizer.Add(self.sort_choice, pos=(4, 2))
        
        left_panel.SetSizer(left_sizer)
        
        # Right panel: Progress, ETA, and log
        right_panel = wx.Panel(panel)
        right_sizer = wx.BoxSizer(wx.VERTICAL)
        
        # Progress bar and ETA
        progress_label = wx.StaticText(right_panel, label="Прогресс:")
        progress_label.SetFont(self.label_font)
        self.progress = wx.Gauge(right_panel, range=100, size=(400, 20))
        self.eta_label = wx.StaticText(right_panel, label="Осталось: 00:00")
        self.eta_label.SetFont(self.label_font)
        
        right_sizer.Add(progress_label, flag=wx.LEFT | wx.BOTTOM, border=5)
        right_sizer.Add(self.progress, flag=wx.EXPAND | wx.LEFT | wx.RIGHT | wx.BOTTOM, border=5)
        right_sizer.Add(self.eta_label, flag=wx.ALIGN_CENTER | wx.BOTTOM, border=5)
        
        # Log area
        log_label = wx.StaticText(right_panel, label="Лог:")
        log_label.SetFont(self.label_font)
        self.log = wx.TextCtrl(right_panel, style=wx.TE_MULTILINE | wx.TE_READONLY | wx.HSCROLL, size=(400, 150))
        self.log.SetFont(self.log_font)
        self.log.SetToolTip("Журнал операций и ошибок")
        
        right_sizer.Add(log_label, flag=wx.LEFT | wx.BOTTOM, border=5)
        right_sizer.Add(self.log, flag=wx.EXPAND | wx.LEFT | wx.RIGHT | wx.BOTTOM, border=5, proportion=1)
        
        right_panel.SetSizer(right_sizer)
        
        # Combine panels
        main_sizer.Add(left_panel, flag=wx.EXPAND | wx.ALL, border=10)
        main_sizer.Add(right_panel, flag=wx.EXPAND | wx.ALL, border=10, proportion=1)
        
        # Bottom buttons
        button_sizer = wx.BoxSizer(wx.HORIZONTAL)
        self.mapping_btn = wx.Button(panel, label="Сопоставление", size=(150, 40))
        self.mapping_btn.SetBackgroundColour(self.theme["button_bg"])
        self.mapping_btn.SetForegroundColour(wx.Colour(255, 255, 255))
        self.mapping_btn.SetFont(self.button_font)
        self.mapping_btn.Enable(False)
        self.mapping_btn.SetToolTip("Свяжите плейсхолдеры PPTX с колонками Excel")
        self.mapping_btn.Bind(wx.EVT_BUTTON, self.open_mapping_window)
        
        self.generate_btn = wx.Button(panel, label="Запустить генерацию", size=(150, 40))
        self.generate_btn.SetBackgroundColour(wx.Colour(76, 175, 80))
        self.generate_btn.SetForegroundColour(wx.Colour(255, 255, 255))
        self.generate_btn.SetFont(self.button_font)
        self.generate_btn.Enable(False)
        self.generate_btn.SetToolTip("Начать генерацию дипломов")
        self.generate_btn.Bind(wx.EVT_BUTTON, self.start_generation)
        
        self.stop_btn = wx.Button(panel, label="Прервать", size=(150, 40))
        self.stop_btn.SetBackgroundColour(wx.Colour(244, 67, 54))
        self.stop_btn.SetForegroundColour(wx.Colour(255, 255, 255))
        self.stop_btn.SetFont(self.button_font)
        self.stop_btn.Enable(False)
        self.stop_btn.SetToolTip("Остановить генерацию")
        self.stop_btn.Bind(wx.EVT_BUTTON, self.stop_generation)
        
        button_sizer.Add(self.mapping_btn, flag=wx.RIGHT, border=10)
        button_sizer.Add(self.generate_btn, flag=wx.RIGHT, border=10)
        button_sizer.Add(self.stop_btn)
        
        # Main sizer with buttons at bottom
        outer_sizer = wx.BoxSizer(wx.VERTICAL)
        outer_sizer.Add(main_sizer, flag=wx.EXPAND, proportion=1)
        outer_sizer.Add(button_sizer, flag=wx.ALIGN_CENTER | wx.TOP | wx.BOTTOM, border=10)
        panel.SetSizer(outer_sizer)
    
    def check_queues(self, event):
        try:
            while True:
                message = self.log_queue.get_nowait()
                message = str(message)
                wx.CallAfter(self.log.AppendText, f"{datetime.now().strftime('%H:%M:%S')}: {message}\n")
        except queue.Empty:
            pass
        try:
            while True:
                progress = self.progress_queue.get_nowait()
                wx.CallAfter(self.progress.SetValue, int(progress))
        except queue.Empty:
            pass
        try:
            while True:
                eta = self.eta_queue.get_nowait()
                eta = str(eta)
                wx.CallAfter(self.eta_label.SetLabel, f"Осталось: {eta}")
        except queue.Empty:
            pass
    
    def log_message(self, message):
        self.log_queue.put(str(message))
    
    def choose_excel(self, event):
        with wx.FileDialog(self, "Выберите Excel-файл", wildcard="Excel files (*.xlsx)|*.xlsx",
                          style=wx.FD_OPEN | wx.FD_FILE_MUST_EXIST) as fileDialog:
            if fileDialog.ShowModal() == wx.ID_OK:
                path = fileDialog.GetPath()
                self.excel_path = path
                self.excel_path_ctrl.SetValue(path)
                self.excel_name.SetLabel(os.path.basename(path))
                self.log_message(f"Загружен Excel: {os.path.basename(path)}")
                self.update_buttons()
    
    def choose_pptx(self, event):
        with wx.FileDialog(self, "Выберите шаблон PPTX", wildcard="PowerPoint files (*.pptx)|*.pptx",
                          style=wx.FD_OPEN | wx.FD_FILE_MUST_EXIST) as fileDialog:
            if fileDialog.ShowModal() == wx.ID_OK:
                path = fileDialog.GetPath()
                self.pptx_path = path
                self.pptx_path_ctrl.SetValue(path)
                self.pptx_name.SetLabel(os.path.basename(path))
                self.log_message(f"Загружен шаблон: {os.path.basename(path)}")
                self.scan_placeholders()
                self.update_buttons()
    
    def choose_output(self, event):
        with wx.DirDialog(self, "Выберите папку для сохранения", style=wx.DD_DEFAULT_STYLE) as dirDialog:
            if dirDialog.ShowModal() == wx.ID_OK:
                path = dirDialog.GetPath()
                self.output_dir = path
                self.output_path_ctrl.SetValue(path)
                self.output_name.SetLabel(os.path.basename(path) or "Папка")
                self.log_message(f"Выбрана папка: {path}")
                self.update_buttons()
    
    def scan_placeholders(self):
        try:
            prs = Presentation(self.pptx_path)
            slide = prs.slides[0]
            self.placeholders = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        matches = re.findall(r"\{([^}]+)\}", paragraph.text)
                        self.placeholders.extend(matches)
            self.placeholders = list(set(self.placeholders))
            self.log_message(f"Найдены плейсхолдеры: {', '.join(self.placeholders)}")
            self.update_sort_choice()
            self.sort_check.Enable(True)
            self.sort_choice.Enable(True)
        except Exception as e:
            self.log_message(f"Ошибка сканирования шаблона: {e}")
            self.sort_check.Enable(False)
            self.sort_choice.Enable(False)
            self.sort_choice.Set([])
    
    def update_sort_choice(self):
        if self.placeholders:
            self.sort_choice.Set(self.placeholders)
            self.sort_column = self.placeholders[0] if self.placeholders else ""
            self.sort_choice.SetSelection(0)
        else:
            self.sort_choice.Set([])
            self.sort_column = ""
    
    def update_buttons(self):
        if self.excel_path and self.pptx_path and self.output_dir and self.placeholders:
            self.mapping_btn.Enable(True)
            if self.column_mapping:
                self.generate_btn.Enable(True)
        else:
            self.mapping_btn.Enable(False)
            self.generate_btn.Enable(False)
    
    def update_error_handling(self, event):
        self.error_handling = {"Остановить": "stop", "Пропустить": "skip", "Заполнить по умолчанию": "default"}[self.error_handling_choice.GetStringSelection()]
        if self.error_handling == "default":
            self.open_default_values_window()
        self.log_message(f"Обработка ошибок: {self.error_handling_choice.GetStringSelection()}")
    
    def update_sorting(self, event):
        self.enable_sorting = self.sort_check.GetValue()
        self.sort_column = self.sort_choice.GetStringSelection()
        self.log_message(f"Сортировка: {'включена' if self.enable_sorting else 'выключена'}, столбец: {self.sort_column}")
    
    def open_default_values_window(self):
        dialog = wx.Dialog(self, title="Значения по умолчанию", size=(400, 300))
        panel = wx.Panel(dialog)
        sizer = wx.BoxSizer(wx.VERTICAL)
        
        self.default_entries = {}
        for placeholder in self.placeholders:
            hsizer = wx.BoxSizer(wx.HORIZONTAL)
            label = wx.StaticText(panel, label=f"{placeholder}:")
            label.SetFont(self.label_font)
            hsizer.Add(label, flag=wx.RIGHT | wx.ALIGN_CENTER_VERTICAL, border=5)
            entry = wx.TextCtrl(panel, value=self.default_values.get(placeholder, "Не указано"))
            hsizer.Add(entry, flag=wx.EXPAND)
            self.default_entries[placeholder] = entry
            sizer.Add(hsizer, flag=wx.EXPAND | wx.ALL, border=5)
        
        save_btn = wx.Button(panel, label="Сохранить")
        save_btn.Bind(wx.EVT_BUTTON, lambda evt: self.save_default_values(dialog))
        sizer.Add(save_btn, flag=wx.ALIGN_CENTER | wx.ALL, border=10)
        
        panel.SetSizer(sizer)
        dialog.ShowModal()
    
    def save_default_values(self, dialog):
        self.default_values = {ph: entry.GetValue() for ph, entry in self.default_entries.items()}
        self.log_message("Значения по умолчанию сохранены")
        dialog.Destroy()
    
    def open_mapping_window(self, event):
        if not self.excel_path or not self.pptx_path:
            wx.MessageBox("Сначала выберите Excel и PPTX!", "Ошибка", wx.OK | wx.ICON_ERROR)
            return
        
        wb = load_workbook(self.excel_path)
        ws = wb.active
        headers = [cell.value or f"Столбец {chr(65+i)}" for i, cell in enumerate(ws[1])]
        
        dialog = wx.Dialog(self, title="Сопоставление плейсхолдеров", size=(400, 400))
        panel = wx.Panel(dialog)
        sizer = wx.BoxSizer(wx.VERTICAL)
        
        self.mapping_choices = {}
        for placeholder in self.placeholders:
            hsizer = wx.BoxSizer(wx.HORIZONTAL)
            label = wx.StaticText(panel, label=f"{placeholder}:")
            label.SetFont(self.label_font)
            hsizer.Add(label, flag=wx.RIGHT | wx.ALIGN_CENTER_VERTICAL, border=5)
            choice = wx.Choice(panel, choices=["Игнорировать"] + headers)
            choice.SetStringSelection(self.column_mapping.get(placeholder, "Игнорировать"))
            hsizer.Add(choice, flag=wx.EXPAND)
            self.mapping_choices[placeholder] = choice
            sizer.Add(hsizer, flag=wx.EXPAND | wx.ALL, border=5)
        
        save_btn = wx.Button(panel, label="Сохранить")
        save_btn.Bind(wx.EVT_BUTTON, lambda evt: self.save_mapping(dialog))
        auto_map_btn = wx.Button(panel, label="Автосопоставление")
        auto_map_btn.Bind(wx.EVT_BUTTON, lambda evt: self.auto_map(headers))
        check_btn = wx.Button(panel, label="Проверить данные")
        check_btn.Bind(wx.EVT_BUTTON, lambda evt: self.check_data(headers, ws))
        
        sizer.Add(save_btn, flag=wx.ALIGN_CENTER | wx.ALL, border=5)
        sizer.Add(auto_map_btn, flag=wx.ALIGN_CENTER | wx.ALL, border=5)
        sizer.Add(check_btn, flag=wx.ALIGN_CENTER | wx.ALL, border=5)
        
        panel.SetSizer(sizer)
        dialog.ShowModal()
    
    def save_mapping(self, dialog):
        self.column_mapping = {
            ph: choice.GetStringSelection()
            for ph, choice in self.mapping_choices.items()
            if choice.GetStringSelection() != "Игнорировать"
        }
        self.log_message("Сопоставление сохранено")
        self.update_buttons()
        dialog.Destroy()
    
    def auto_map(self, headers):
        synonyms = {
            "NAME": ["ФИО", "Имя", "Name", "Full Name"],
            "REGN": ["Номер", "Регистрация", "Reg", "ID"],
            "LEARN": ["Курс", "Программа", "Learn", "Course"],
            "TIME": ["Часы", "Время", "Hours", "Duration"],
            "DATE": ["Дата", "Date"]
        }
        for placeholder in self.placeholders:
            choice = self.mapping_choices[placeholder]
            for header in headers:
                if header and any(syn.lower() in header.lower() for syn in synonyms.get(placeholder, [placeholder])):
                    choice.SetStringSelection(header)
                    break
        self.log_message("Выполнено автосопоставление")
    
    def check_data(self, headers, ws):
        errors = []
        for row_idx, row in enumerate(ws.iter_rows(min_row=2, values_only=True), 2):
            for placeholder, col_name in self.column_mapping.items():
                col_idx = headers.index(col_name) if col_name in headers else ord(col_name) - ord('A')
                if row[col_idx] is None:
                    errors.append(f"Строка {row_idx}: пустое поле {placeholder} ({col_name})")
        if errors:
            wx.MessageBox(
                "\n".join(errors[:5]) + (f"\n...и ещё {len(errors)-5} ошибок" if len(errors) > 5 else ""),
                "Предупреждение", wx.OK | wx.ICON_WARNING
            )
        else:
            wx.MessageBox("Ошибок не найдено", "Проверка", wx.OK | wx.ICON_INFORMATION)
    
    def start_generation(self, event):
        if not all(self.column_mapping.values()):
            wx.MessageBox("Не все плейсхолдеры сопоставлены!", "Ошибка", wx.OK | wx.ICON_ERROR)
            return
        self.generate_btn.Enable(False)
        self.stop_btn.Enable(True)
        self.progress.SetValue(0)
        self.eta_label.SetLabel("Осталось: 00:00")
        self.stop_event.clear()
        self.generation_thread = threading.Thread(target=self.run_generation)
        self.generation_thread.start()
    
    def run_generation(self):
        try:
            success = generate_diplomas(
                self.excel_path, self.pptx_path, self.output_dir,
                self.column_mapping, self.error_handling, self.default_values,
                self.font_settings, self.sort_column, self.enable_sorting,
                self.log_queue, self.progress_queue, self.eta_queue, self.stop_event
            )
            if success:
                wx.CallAfter(wx.MessageBox, f"Дипломы сгенерированы в: {self.output_dir}", "Успех", wx.OK | wx.ICON_INFORMATION)
        except Exception as e:
            wx.CallAfter(wx.MessageBox, str(e), "Ошибка", wx.OK | wx.ICON_ERROR)
        finally:
            wx.CallAfter(self.reset_buttons)
    
    def stop_generation(self, event):
        self.stop_event.set()
        self.progress.SetValue(0)
        self.log_message("Прерывание генерации...")
        threading.Thread(target=self.cleanup_powerpoint, daemon=True).start()
        if self.generation_thread:
            self.generation_thread.join(timeout=5)
        self.reset_buttons()
    
    def cleanup_powerpoint(self):
        for proc in psutil.process_iter(['name']):
            if proc.info['name'].lower() == 'powerpnt.exe':
                proc.terminate()
                try:
                    proc.wait(timeout=3)
                except psutil.TimeoutExpired:
                    proc.kill()
        for file in os.listdir():
            if file.startswith("temp_") and file.endswith((".pptx", ".jpg", ".pdf")):
                try:
                    os.remove(file)
                except:
                    pass
    
    def reset_buttons(self):
        self.generate_btn.Enable(bool(self.column_mapping))
        self.stop_btn.Enable(False)
        self.eta_label.SetLabel("Осталось: 00:00")
    
    def on_closing(self, event):
        if self.generation_thread and self.generation_thread.is_alive():
            self.stop_event.set()
            self.cleanup_powerpoint()
            self.generation_thread.join(timeout=5)
        for file in os.listdir():
            if file.startswith("temp_") and file.endswith((".pptx", ".jpg", ".pdf")):
                try:
                    os.remove(file)
                except:
                    pass
        self.timer.Stop()
        self.Destroy()
    
    def load_config(self):
        try:
            with open("config.json", "r", encoding="utf-8") as f:
                config = json.load(f)
                self.excel_path = config.get("excel_path", "")
                self.pptx_path = config.get("pptx_path", "")
                self.output_dir = config.get("output_dir", "")
                self.column_mapping = config.get("column_mapping", {})
                self.error_handling = config.get("error_handling", "stop")
                self.default_values = config.get("default_values", {})
                self.sort_column = config.get("sort_column", "")
                self.enable_sorting = config.get("enable_sorting", True)
                if self.excel_path:
                    self.excel_path_ctrl.SetValue(self.excel_path)
                    self.excel_name.SetLabel(os.path.basename(self.excel_path))
                if self.pptx_path:
                    self.pptx_path_ctrl.SetValue(self.pptx_path)
                    self.pptx_name.SetLabel(os.path.basename(self.pptx_path))
                    self.scan_placeholders()
                if self.output_dir:
                    self.output_path_ctrl.SetValue(self.output_dir)
                    self.output_name.SetLabel(os.path.basename(self.output_dir) or "Папка")
                self.error_handling_choice.SetStringSelection({"stop": "Остановить", "skip": "Пропустить", "default": "Заполнить по умолчанию"}[self.error_handling])
                self.sort_check.SetValue(self.enable_sorting)
                if self.sort_column and self.placeholders:
                    self.sort_choice.SetSelection(self.placeholders.index(self.sort_column) if self.sort_column in self.placeholders else 0)
                self.update_buttons()
        except FileNotFoundError:
            pass

if __name__ == "__main__":
    app = wx.App()
    frame = DiplomaGeneratorApp()
    frame.Show()
    app.MainLoop()